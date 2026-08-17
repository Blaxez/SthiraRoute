"""
Assembles the rendered slides into a .pptx at the official SIH page size.

The deck is authored in HTML and rendered by Chromium, so the PDF is the
primary deliverable; this produces a PowerPoint file for anyone who wants one.
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Inches

HERE = Path(__file__).parent
SLIDES = sorted((HERE / "out" / "slides").glob("*.png"))
OUT = HERE / "out" / "SIH2026-HackShastra-SthiraRoute.pptx"

if not SLIDES:
    raise SystemExit("no rendered slides — run apps/web/render-deck.mjs first")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

for png in SLIDES:
    slide = prs.slides.add_slide(blank)
    slide.shapes.add_picture(str(png), 0, 0,
                             width=prs.slide_width, height=prs.slide_height)

try:
    prs.save(str(OUT))
except PermissionError:
    # The deck is almost always open in PowerPoint when you re-render it.
    # Write beside it rather than losing the run.
    OUT = OUT.with_stem(OUT.stem + "-new")
    prs.save(str(OUT))
    print("  ! the existing .pptx is open in PowerPoint — wrote a copy instead")

print(f"saved {OUT}")
print(f"  {len(SLIDES)} slides at "
      f"{Emu(prs.slide_width).inches:.3f} x {Emu(prs.slide_height).inches:.3f} in")
for p in SLIDES:
    print(f"  · {p.name}")
